from pptx import Presentation
import os
import shutil

# ------------------- CONFIGURACIÓN DE USUARIO -------------------

# Ruta de entrada: Plantilla PowerPoint
INPUT_PPTX_TEMPLATE = './Inputs/PlantillaPyCoffee.pptx'

# Ruta de salida
OUTPUT_PATH = '../../Outputs'

# Parámetros de las diapositivas
SLIDE_LIST = [
    # Diapositiva Título
    {
        'texto_titulo': "Revisión Anual",
        'layout': 0,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "Edu Morales", 'nivel': 0},
                ]
            }
        ]
    },
    # Diapositiva Título y Objeto - ¿Qué vamos a ver?
    {
        'texto_titulo': "¿Qué vamos a ver?",
        'layout': 1,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "Comparación de resultados entre el año 2023 y 2024", 'nivel': 0},
                    {'texto': "Previsiones para el año 2025", 'nivel': 0},
                ]
            }
        ]
    },
    # Diapositiva Encabezado de sección - Comparación de resultados
    {
        'texto_titulo': "Comparación de resultados",
        'layout': 2,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "2023 - 2024", 'nivel': 0},
                ]
            }
        ]
    },
    # Diapositiva Comparación - Comparación de resultados
    {
        'texto_titulo': "Comparación de resultados",
        'layout': 4,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "2023", 'nivel': 0},
                ]
            },
            {
                'placeholder_id': 2,
                'elementos': [
                    {'texto': "Venta de 3.231 PyCoffeeMax R340", 'nivel': 0},
                    {'texto': "Venta de 1.503 PyCoffeeMax R570", 'nivel': 0},
                    {'texto': "Venta de 741 PyCoffeeMax X3000", 'nivel': 0},
                    {'texto': "Facturación: 99.920€", 'nivel': 0},
                ]
            },
            {
                'placeholder_id': 3,
                'elementos': [
                    {'texto': "2024", 'nivel': 0},
                ]
            },
            {
                'placeholder_id': 4,
                'elementos': [
                    {'texto': "Venta de 6.134 PyCoffeeMax R340", 'nivel': 0},
                    {'texto': "Venta de 1.398 PyCoffeeMax R570", 'nivel': 0},
                    {'texto': "Venta de 3.150 PyCoffeeMax X3000", 'nivel': 0},
                    {'texto': "Facturación: 246.800€", 'nivel': 0},
                ]
            },
        ]
    },
    # Diapositiva Contenido con Título - Previsiones 2025
    {
        'texto_titulo': "Previsiones 2025",
        'layout': 7,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "PyCoffeeMax R340: + 37%", 'nivel': 0},
                    {'texto': "PyCoffeeMax R570: + 42%", 'nivel': 0},
                    {'texto': "PyCoffeeMax X3000: + 15%", 'nivel': 0},
                    {'texto': "Facturación estimada: 304.910€", 'nivel': 0},
                ]
            },
            {
                'placeholder_id': 2,
                'elementos': [
                    {'texto': "PyCoffee", 'nivel': 0},
                ]
            }
        ]
    },
    # Diapositiva Agradecimiento - Muchas gracias
    {
        'texto_titulo': "Muchas gracias",
        'layout': 11,
        'placeholders': [
        ]
    }
]

# ------------------- /CONFIGURACIÓN DE USUARIO ------------------

# Eliminar y crear carpeta de nuevo
def RestaurarCarpeta(folder_path):
    # Verificar si la carpeta existe
    if os.path.exists(folder_path):
        # Eliminar la carpeta
        shutil.rmtree(folder_path)

    # Crear carpeta
    os.makedirs(folder_path)


# Crear diapositiva
def CrearDiapositiva(presentacion, layout_id):
    # Crear diapositiva
    slide_layout = presentacion.slide_layouts[layout_id]
    slide = presentacion.slides.add_slide(slide_layout)

    return slide


# Añadir texto a un placeholder de una diapositiva
def AñadirTextoPlaceholder(slide, placeholder_id, texto, nivel=0):
    # Crear el cuadro de texto (textframe)
    body_shape = slide.shapes.placeholders[placeholder_id]
    tf = body_shape.text_frame

    # Comprobbar si ya hay texto en el text_frame
    if tf.text:
        # Añadir un nuevo párrafo
        p = tf.add_paragraph()
        p.text = texto
        p.level = nivel
    else:
        # Reemplazar el texto existente
        tf.text = texto
        tf.level = nivel


# Rutina principal
def main():
    # Eliminar y crear OUTPUT_PATH de nuevo
    RestaurarCarpeta(OUTPUT_PATH)

    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Iteración sobre la lista de diccionarios SLIDE_LIST
    for slide_dict in SLIDE_LIST:
        # Crear diapositiva
        layout_id = slide_dict['layout']
        slide = CrearDiapositiva(prs, layout_id)

        # Modificar el título de la diapositiva (placeholder 0)
        texto = slide_dict['texto_titulo']
        AñadirTextoPlaceholder(slide, 0, texto)

        # Iterar sobre los placeholders
        for ph_dict in slide_dict['placeholders']:
            placeholder_id = ph_dict['placeholder_id']

            # Iterar sobre cada elemento del placeholder
            for elem_dict in ph_dict['elementos']:
                # Añadir texto al placeholder
                nivel = elem_dict['nivel']
                texto = elem_dict['texto']
                AñadirTextoPlaceholder(slide, placeholder_id, texto, nivel)

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "PlantillasProyecto_Output.pptx")

    # Mensaje de presentación creada
    print("")
    print("La presentación se ha creado correctamente!")


if __name__ == '__main__':
    main()

