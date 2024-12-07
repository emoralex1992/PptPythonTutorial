from pptx import Presentation
import os
import shutil
import pandas as pd
import warnings
from pptx.util import Pt, Cm
import openpyxl
from pptx.dml.color import RGBColor
warnings.filterwarnings("ignore", category=UserWarning, module='openpyxl')

# ------------------- CONFIGURACIÓN DE USUARIO -------------------

# Ruta de entrada: Plantilla PowerPoint
INPUT_PPTX_TEMPLATE = './Inputs/Template_GCU.pptx'

# Ruta de entrada: Fichero Excel de configuración de la presentación
INPUT_EXCEL_CFG = './Inputs/Presentacion_Cfg.xlsx'

# Ruta de salida
OUTPUT_PATH = './Outputs'

# ------------------- /CONFIGURACIÓN DE USUARIO ------------------

# Eliminar y crear carpeta de nuevo
def RestaurarCarpeta(folder_path):
    # Verificar si la carpeta existe
    if os.path.exists(folder_path):
        # Eliminar la carpeta
        shutil.rmtree(folder_path)

    # Crear carpeta
    os.makedirs(folder_path)


# Crear diapositiva
def CrearDiapositiva(presentacion, layout_id):
    # Crear diapositiva
    slide_layout = presentacion.slide_layouts[layout_id]
    slide = presentacion.slides.add_slide(slide_layout)

    return slide


# Añadir texto a un placeholder de una diapositiva
def AñadirTextoPlaceholder(slide, placeholder_id, runs_list, nivel=0, paragraph_idx=0):
    # Crear el cuadro de texto (textframe)
    body_shape = slide.shapes.placeholders[placeholder_id]
    tf = body_shape.text_frame

    # Comprobar si ya hay texto en el text_frame
    if paragraph_idx == 0:
        # Utilizar el párrafo 0
        p = tf.paragraphs[0]
    else:
        # Añadir un nuevo párrafo
        p = tf.add_paragraph()

    # Asignar nivel
    p.level = nivel

    # Iterar sobre runs_list
    for run_dict in runs_list:
        AñadirRunAlParrafo(p, run_dict)


# Rutina para añadir runs al párrafo
def AñadirRunAlParrafo(paragraph, run_dict):
    run = paragraph.add_run()
    run.text = run_dict['texto']

    # Añadir negrita
    if "negrita" in run_dict:
        run.font.bold = run_dict["negrita"]

    # Añadir cursiva
    if "cursiva" in run_dict:
        run.font.italic = run_dict["cursiva"]

    # Añadir subrayado
    if "subrayado" in run_dict:
        run.font.underline = run_dict["subrayado"]

    # Añadir tamaño del texto
    if "tamaño_txt" in run_dict and run_dict['tamaño_txt'] != None:
        run.font.size = Pt(run_dict["tamaño_txt"])

    # Añadir color de fuente
    if "color_txt" in run_dict and run_dict['color_txt'] != None:
        run.font.color.rgb = RGBColor(*run_dict["color_txt"])

    # Añadir tipo de fuente
    if "fuente" in run_dict and run_dict['fuente'] != None:
        run.font.name = run_dict["fuente"]


# Rutina para añadir una imagen a un placeholder
def AñadirImagenPlaceholder(slide, img_path, placeholder, img_param_dict):
    # Insertar imagen en la diapositiva
    img = slide.shapes.add_picture(img_path, placeholder.left, placeholder.top)

    # Calcular la relación de aspecto del marcador de posición
    placeholder_aspect_ratio = placeholder.width / placeholder.height

    # Calcular la relación de aspecto de la imagen
    img_aspect_ratio = img.width / img.height

    # Recortar imagen de forma automática para ajustarla al placeholder
    if img_param_dict['img_modo'] == "RECORTAR PARA ENCAJAR":
        # Ajustar el ancho y alto de la imagen al ancho y alto del placeholder
        img.height = placeholder.height
        img.width = placeholder.width

        # La imagen es más ancha que el placeholder
        if img_aspect_ratio > placeholder_aspect_ratio:
            # Calcular la proporción de anchura
            image_prop_width = int(img.height * img_aspect_ratio)

            # Recortar la imagen para ajustarla a la anchura del placeholder basándonos en la anchura proporcional
            crop_value = ((image_prop_width - img.width) / image_prop_width) / 2
            img.crop_left = crop_value
            img.crop_right = crop_value

        # La imagen es más alta que el placeholder
        else:
            # Calcular la proporción de altura
            image_prop_height = int(img.width / img_aspect_ratio)

            # Recortar la imagen para ajustarla a la altura del placeholder basándonos en la altura proporcional
            crop_value = ((image_prop_height - img.height) / image_prop_height) / 2
            img.crop_top = crop_value
            img.crop_bottom = crop_value

        # Aplicar offset al recorte
        img.crop_left += img_param_dict['rec_offset_der_porc'] - img_param_dict['rec_offset_izq_porc']
        img.crop_right += img_param_dict['rec_offset_izq_porc'] - img_param_dict['rec_offset_der_porc']
        img.crop_top += img_param_dict['rec_offset_inf_porc'] - img_param_dict['rec_offset_sup_porc']
        img.crop_bottom += img_param_dict['rec_offset_sup_porc'] - img_param_dict['rec_offset_inf_porc']

    # Imagen con ancho y alto automáticos --> Ajustar al placeholder
    elif img_param_dict['img_modo'] == "ANCHO/ALTO AUTOMATICO":
        # La imagen es más ancha que el placeholder
        if img_aspect_ratio > placeholder_aspect_ratio:
            scale_factor = placeholder.width / img.width
        # La imagen es más alta que el placeholder
        else:
            scale_factor = placeholder.height / img.height

        # Aplicar el factor de escala a la imagen
        img.width = int(img.width * scale_factor)
        img.height = int(img.height * scale_factor)

    # Imagen con ancho y alto personalizados
    elif img_param_dict['img_modo'] == "TAMAÑO PERSONALIZADO":
        # Ancho y Alto definidos
        if (img_param_dict['ancho_img_cm'] != None) and (img_param_dict['alto_img_cm'] != None):
            img.width = Cm(img_param_dict['ancho_img_cm'])
            img.height = Cm(img_param_dict['alto_img_cm'])

        # Únicamente el Ancho de la imagen ha sido definido
        elif (img_param_dict['ancho_img_cm'] != None) and (img_param_dict['alto_img_cm'] == None):
            img.height = Cm(img_param_dict['ancho_img_cm'] * (img.height / img.width))
            img.width = Cm(img_param_dict['ancho_img_cm'])

        # Únicamente el Alto de la imagen ha sido definido
        elif (img_param_dict['ancho_img_cm'] == None) and (img_param_dict['alto_img_cm'] != None):
            img.width = Cm(img_param_dict['alto_img_cm'] * (img.width / img.height))
            img.heigh = Cm(img_param_dict['alto_img_cm'])

        # Error: Ni el Alto ni el Ancho de la imagen han sido definidos
        else:
            raise ValueError(f"El modo de imagen TAMAÑO PERSONALIZADO debe tener uno o todos los campos \"ALTO_IMG_CM\" y/o \"ANCHO_IMG_CM\" definidos "
                             f"para la imagen: {img_path}")

    # Alinear la imagen verticalmente dentro del placeholder ("IZQUIERDA" por defecto)
    if img_param_dict['img_alinear_v'] == "DERECHA":
        width_diff = placeholder.width - img.width
        img.left = placeholder.left + width_diff
    elif img_param_dict['img_alinear_v'] == "CENTRO":
        width_diff = placeholder.width - img.width
        img.left = placeholder.left + int(width_diff / 2)

    # Alinear la imagen horizontalmente dentro del placeholder ("ARRIBA" por defecto)
    if img_param_dict['img_alinear_h'] == "ABAJO":
        height_diff = placeholder.height - img.height
        img.top = placeholder.top + height_diff
    elif img_param_dict['img_alinear_h'] == "CENTRO":
        height_diff = placeholder.height - img.height
        img.top = placeholder.top + int(height_diff / 2)

# Imprimir todos los nombres de los diferentes diseños de diapositivas (layouts) de una plantilla PowerPoint
def ImprimirDiseñosPresentacion(prs):
    print("LISTA DE DISEÑOS DE DIAPOSITIVAS:")
    # Iterar sobre los layouts
    for layout in prs.slide_layouts:
        # Imprimir el nombre del layout
        print(layout.name)
    print("")


# Crear presentación con cada uno de los diseños y el nombre de los placeholders
def ExportarNombrePlaceholders():
    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Iterar sobre los layouts
    for layout_id, layout in enumerate(prs.slide_layouts):
        slide = CrearDiapositiva(prs, layout_id)
        # print(layout.name)

        # Iterar sobre los placeholders
        for ph in layout.placeholders:
            try:
                text = str(ph.placeholder_format.idx) + " - " + ph.name
                runs_list = [
                    {
                        'texto': text,
                    }
                ]
                AñadirTextoPlaceholder(slide, ph.placeholder_format.idx, runs_list, 0, 0)
                # print(f"   {ph.placeholder_format.idx} - {ph.name}")
            except KeyError:
                i = 0

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "Marcadores.pptx")

# Convertir color hexadecimal a RGB
def HexToRgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return [int(hex_color[i:i+2], 16) for i in (0, 2, 4)]

# Importar datos desde el fichero de configuración Excel
def ImportarDatosExcel(prs):
    # Inicializar lista slide_list
    slide_list = []

    # Leer fichero Excel
    df = pd.read_excel(INPUT_EXCEL_CFG, sheet_name='Presentación')

    # Leer fichero Excel utilizando openpyxl
    wb = openpyxl.load_workbook(INPUT_EXCEL_CFG)
    ws = wb['Presentación']

    # Iterar sobre las filas del DataFrame
    for index, row in df.iterrows():
        # Leer celdas de la fila
        nombre_layout = row['DISEÑO_DIAPOSITIVA']
        placeholder_id = row['NUM_MARCADOR']
        nivel = int(row['NIVEL_TEXTO'].replace("N", "")) if pd.notna(row['NIVEL_TEXTO']) else None
        texto_run = str(row['TEXTO']) if pd.notna(row['TEXTO']) else ""
        negrita = row['NEGRITA'] == "Sí"
        cursiva = row['CURSIVA'] == "Sí"
        subrayado = row['SUBRAYADO'] == "Sí"
        tamaño_texto = int(row['TAMAÑO_TXT']) if pd.notna(row['TAMAÑO_TXT']) else None
        fuente = row['FUENTE'] if pd.notna(row['FUENTE']) else None
        img_path = row['IMG_PATH'] if pd.notna(row['IMG_PATH']) else None
        img_modo = row['IMG_MODO'] if pd.notna(row['IMG_MODO']) else None
        rec_offset_izq_porc = float(row['REC_OFFSET_IZQ_PORC']) if pd.notna(row['REC_OFFSET_IZQ_PORC']) else 0.0
        rec_offset_der_porc = float(row['REC_OFFSET_DER_PORC']) if pd.notna(row['REC_OFFSET_DER_PORC']) else 0.0
        rec_offset_sup_porc = float(row['REC_OFFSET_SUP_PORC']) if pd.notna(row['REC_OFFSET_SUP_PORC']) else 0.0
        rec_offset_inf_porc = float(row['REC_OFFSET_INF_PORC']) if pd.notna(row['REC_OFFSET_INF_PORC']) else 0.0
        ancho_img_cm = float(row['ANCHO_IMG_CM']) if pd.notna(row['ANCHO_IMG_CM']) else None
        alto_img_cm = float(row['ALTO_IMG_CM']) if pd.notna(row['ALTO_IMG_CM']) else None
        img_alinear_v = row['IMG_ALINEAR_V'] if pd.notna(row['IMG_ALINEAR_V']) else None
        img_alinear_h = row['IMG_ALINEAR_H'] if pd.notna(row['IMG_ALINEAR_H']) else None

        # Obtener el color de la celda correspondiente a COLOR_TXT
        cell = ws.cell(row=index + 2, column=df.columns.get_loc('COLOR_TXT') + 1)
        color_celda = cell.fill.fgColor
        r, g, b = None, None, None  # Inicializar valores de color

        if color_celda.type == "rgb" and isinstance(color_celda.rgb, str):  # Si es un color directo en formato hexadecimal
            color_hex = color_celda.rgb
            if len(color_hex) == 8 and color_hex != "00000000":  # Verificar si no es transparente o sin color asignado
                r, g, b = HexToRgb(color_hex[2:])  # Convertir a RGB

        elif color_celda.type == "theme":  # Si es un color de tema
            # Lanzar un error informando al usuario que debe cambiar el color de tema
            raise ValueError(f"El color en la fila {index + 2} columna 'COLOR_TXT' es un color de tema. "
                             f"Por favor, cambie el color a un valor explícito (no de tema) en el archivo Excel.")

        # Diapositiva encontrada
        if pd.notna(nombre_layout):
            # Encontrar el layout_id correspondiente al layout_name
            layout_id = None
            for i, layout in enumerate(prs.slide_layouts):
                if layout.name == nombre_layout:
                    layout_id = i
                    break

            # Error: layout_id no encontrado
            if layout_id is None:
                row_error = 2 + index
                raise ValueError(f"No se encontró un diseño de diapositiva con el nombre: {nombre_layout} - FILA {row_error} / COLUMNA \"DISEÑO_DIAPOSITIVA\"")

            print(f"{nombre_layout} - Idx: {layout_id}")

            # Crear diccionario para la diapositiva y añadirlo a slide_list
            slide_list.append({
                'layout_name': nombre_layout,
                'layout': layout_id,
                'placeholders': []
            })

        # Marcador (placeholder) encontrado
        elif pd.notna(placeholder_id):
            # Verificar si el placeholder_id es válido para el diseño de la diapositiva actual
            layout_id = slide_list[-1]['layout']
            valid_placeholder_ids = [ph.placeholder_format.idx for ph in prs.slide_layouts[layout_id].placeholders]

            # Error: el placeholder no se encuentra en el diseño de diapositiva
            if placeholder_id not in valid_placeholder_ids:
                row_error = 2 + index
                raise ValueError(f"No se encontró el marcaddor con ID {int(placeholder_id)} en el diseño de diapositiva \"{slide_list[-1]['layout_name']}\" - FILA {row_error} / COLUMNA \"NUM_MARCADOR\"")

            # Añadir el placeholder al último slide_dict en slide_list
            slide_list[-1]['placeholders'].append({
                'placeholder_id': placeholder_id,
                'elementos': []
            })

        # Nivel_Texto encontrado
        elif nivel != None:
            # Añadir el texto a los elementos del último placeholder
            slide_list[-1]['placeholders'][-1]['elementos'].append({
                'runs': [],
                'nivel': nivel
            })

        # Texto encontrado
        elif texto_run != "":
            # Añadir run a lista "runs"
            slide_list[-1]['placeholders'][-1]['elementos'][-1]['runs'].append(
                {
                    'texto': texto_run,
                    'negrita': negrita,
                    'cursiva': cursiva,
                    'subrayado': subrayado,
                    'tamaño_txt': tamaño_texto,
                    'color_txt': [r, g, b] if r != None else None,
                    'fuente': fuente
                })

        # Imagen encontrada
        elif img_path != None:
            # Añadir el texto a los elementos del último placeholder
            slide_list[-1]['placeholders'][-1]['elementos'].append({
                'img_path': img_path,
                'img_param_dict': {
                    'img_modo': img_modo,
                    'rec_offset_izq_porc': rec_offset_izq_porc,
                    'rec_offset_der_porc': rec_offset_der_porc,
                    'rec_offset_sup_porc': rec_offset_sup_porc,
                    'rec_offset_inf_porc': rec_offset_inf_porc,
                    'ancho_img_cm': ancho_img_cm,
                    'alto_img_cm': alto_img_cm,
                    'img_alinear_v': img_alinear_v,
                    'img_alinear_h': img_alinear_h,
                }
            })

    return slide_list


# Rutina principal
def main():
    # Eliminar y crear OUTPUT_PATH de nuevo
    RestaurarCarpeta(OUTPUT_PATH)

    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Imprimir los diseños de la presentación por consola
    ImprimirDiseñosPresentacion(prs)

    # Exportar el nombre de los placeholders en una nueva presentación
    ExportarNombrePlaceholders()

    # Importar datos desde fichero Excel
    slide_list = ImportarDatosExcel(prs)

    # Iteración sobre la lista de diccionarios SLIDE_LIST
    for slide_dict in slide_list:
        # Crear diapositiva
        layout_id = slide_dict['layout']
        slide = CrearDiapositiva(prs, layout_id)

        # Iterar sobre los placeholders
        for ph_dict in slide_dict['placeholders']:
            placeholder_id = ph_dict['placeholder_id']

            # Iterar sobre cada elemento del placeholder
            for elem_idx, elem_dict in enumerate(ph_dict['elementos']):
                # Texto detectado
                if ('nivel' in elem_dict):
                    # Añadir texto al placeholder
                    nivel = elem_dict['nivel']
                    runs_list = elem_dict['runs']
                    AñadirTextoPlaceholder(slide, placeholder_id, runs_list, nivel, elem_idx)

                # Imagen detectada
                elif ('img_path' in elem_dict):
                    ph = slide.placeholders[placeholder_id]
                    AñadirImagenPlaceholder(slide, elem_dict['img_path'], ph, elem_dict['img_param_dict'])

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "PlantillasProyecto_Output.pptx")

    # Mensaje de presentación creada
    print("")
    print("La presentación se ha creado correctamente!")


if __name__ == '__main__':
    main()

