from pptx import Presentation
from pptx.dml.color import RGBColor
import os
import shutil

# ------------------- CONFIGURACIÓN DE USUARIO -------------------

# Ruta de salida
OUTPUT_PATH = '../../Outputs'

# Parámetros de las diapositivas
SLIDE_LIST = [
    # Diapositiva Amarilla
    {
        'texto_titulo': "Título Diapositiva Amarilla",
        'color_fondo_rgb': [255, 255, 0],    # Amarillo estándar
        'layout': 0,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "Elemento nivel 0 - Amarillo", 'nivel': 0},
                    {'texto': "Elemento nivel 4 - Amarillo", 'nivel': 4},
                    {'texto': "Elemento nivel 2 - Amarillo", 'nivel': 2},
                ]
            }
        ]
    },
    # Diapositiva Verde
    {
        'texto_titulo': "Título Diapositiva Verde",
        'color_fondo_rgb': [0, 255, 0],        # Verde estándar
        'layout': 2,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "Elemento nivel 1 - Verde", 'nivel': 1},
                    {'texto': "Elemento nivel 2 - Verde", 'nivel': 2},
                    {'texto': "Elemento nivel 5 - Verde", 'nivel': 5},
                ]
            }
        ]
    },
    # Diapositiva Naranja
    {
        'texto_titulo': "Título Diapositiva Naranja",
        'color_fondo_rgb': [255, 153, 51],     # Naranja estándar
        'layout': 9,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "Elemento nivel 2 - Naranja", 'nivel': 2},
                    {'texto': "Elemento nivel 3 - Naranja", 'nivel': 3},
                    {'texto': "Elemento nivel 4 - Naranja", 'nivel': 4},
                ]
            }
        ]
    },
    # Diapositiva Fucsia
    {
        'texto_titulo': "Título Diapositiva Fucsia",
        'color_fondo_rgb': [255, 0, 255],      # Fucsia estándar
        'layout': 8,
        'placeholders': [
            {
                'placeholder_id': 1,
                'elementos': [
                    {'texto': "Elemento nivel 0 - Fucsia", 'nivel': 0},
                    {'texto': "Elemento nivel 1 - Fucsia", 'nivel': 1},
                    {'texto': "Elemento nivel 2 - Fucsia", 'nivel': 2},
                ]
            },
            {
                'placeholder_id': 2,
                'elementos': [
                    {'texto': "Subtítulo nivel 0 - Fucsia", 'nivel': 0},
                    {'texto': "Subtítulo nivel 1 - Fucsia", 'nivel': 1},
                    {'texto': "Subtítulo nivel 2 - Fucsia", 'nivel': 2},
                ]
            }
        ]
    },
]

# ------------------- /CONFIGURACIÓN DE USUARIO ------------------

# Eliminar y crear carpeta de nuevo
def RestaurarCarpeta(folder_path):
    # Verificar si la carpeta existe
    if os.path.exists(folder_path):
        # Eliminar la carpeta
        shutil.rmtree(folder_path)

    # Crear carpeta
    os.makedirs(folder_path)


# Crear diapositiva
def CrearDiapositiva(presentacion, layout_id):
    # Crear diapositiva
    slide_layout = presentacion.slide_layouts[layout_id]
    slide = presentacion.slides.add_slide(slide_layout)

    return slide


# Añadir texto a un placeholder de una diapositiva
def AñadirTextoPlaceholder(slide, placeholder_id, texto, nivel=0):
    # Crear el cuadro de texto (textframe)
    body_shape = slide.shapes.placeholders[placeholder_id]
    tf = body_shape.text_frame

    # Comprobbar si ya hay texto en el text_frame
    if tf.text:
        # Añadir un nuevo párrafo
        p = tf.add_paragraph()
        p.text = texto
        p.level = nivel
    else:
        # Reemplazar el texto existente
        tf.text = texto
        tf.level = nivel


# Cambiar color del fondo de la diapositiva
def CambiarColorFondo(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


# Rutina principal
def main():
    # Eliminar y crear OUTPUT_PATH de nuevo
    RestaurarCarpeta(OUTPUT_PATH)

    # Crear Presentación
    prs = Presentation()

    # Iteración sobre la lista de diccionarios SLIDE_LIST
    for slide_dict in SLIDE_LIST:
        # Crear diapositiva
        layout_id = slide_dict['layout']
        slide = CrearDiapositiva(prs, layout_id)

        # Cambiar fondo de la diapositiva
        r = slide_dict['color_fondo_rgb'][0]
        g = slide_dict['color_fondo_rgb'][1]
        b = slide_dict['color_fondo_rgb'][2]
        CambiarColorFondo(slide, r, g, b)
        
        # Modificar el título de la diapositiva (placeholder 0)
        texto = slide_dict['texto_titulo']
        AñadirTextoPlaceholder(slide, 0, texto)

        # Iterar sobre los placeholders
        for ph_dict in slide_dict['placeholders']:
            placeholder_id = ph_dict['placeholder_id']

            # Iterar sobre cada elemento del placeholder
            for elem_dict in ph_dict['elementos']:
                # Añadir texto al placeholder
                nivel = elem_dict['nivel']
                texto = elem_dict['texto']
                AñadirTextoPlaceholder(slide, placeholder_id, texto, nivel)


    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "1_DiapositivasProyecto.pptx")

    # Mensaje de presentación creada
    print("")
    print("La presentación se ha creado correctamente!")


if __name__ == '__main__':
    main()

