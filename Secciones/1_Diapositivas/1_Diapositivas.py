from pptx import Presentation
from pptx.dml.color import RGBColor
import os
import shutil

# ------------------- CONFIGURACIÓN DE USUARIO -------------------

# Ruta de salida
OUTPUT_PATH = '../../Outputs'

# ------------------- /CONFIGURACIÓN DE USUARIO ------------------

# Eliminar y crear carpeta de nuevo
def RestaurarCarpeta(folder_path):
    # Verificar si la carpeta existe
    if os.path.exists(folder_path):
        # Eliminar la carpeta
        shutil.rmtree(folder_path)

    # Crear carpeta
    os.makedirs(folder_path)


# Crear diapositiva
def CrearDiapositiva(presentacion, layout_id):
    # Crear diapositiva
    slide_layout = presentacion.slide_layouts[layout_id]
    slide = presentacion.slides.add_slide(slide_layout)

    return slide


# Añadir texto a un placeholder de una diapositiva
def AñadirTextoPlaceholder(slide, placeholder_id, texto, nivel=0):
    # Crear el cuadro de texto (textframe)
    body_shape = slide.shapes.placeholders[placeholder_id]
    tf = body_shape.text_frame

    # Comprobbar si ya hay texto en el text_frame
    if tf.text:
        # Añadir un nuevo párrafo
        p = tf.add_paragraph()
        p.text = texto
        p.level = nivel
    else:
        # Reemplazar el texto existente
        tf.text = texto
        tf.level = nivel


# Cambiar color del fondo de la diapositiva
def CambiarColorFondo(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


# Rutina principal
def main():
    # Eliminar y crear OUTPUT_PATH de nuevo
    RestaurarCarpeta(OUTPUT_PATH)

    # Crear Presentación
    prs = Presentation()

    # Crear diapositiva
    slide = CrearDiapositiva(prs, 1)

    # Añadir texto al placeholder 0 (título)
    AñadirTextoPlaceholder(slide, 0, "Nuevo título")

    # Añadir texto al placeholder 1 --> primer elemento (nivel 0 por defecto)
    AñadirTextoPlaceholder(slide, 1, "Elemento 1")

    # Añadir texto al placeholder 1 --> segundo elemento (nivel 0 por defecto)
    AñadirTextoPlaceholder(slide, 1, "Elemento 2")

    # Añadir texto al placeholder 1 --> segundo elemento (nivel 1)
    AñadirTextoPlaceholder(slide, 1, "Subelemento 2.1", 1)

    # Cambiar fondo de la diapositiva
    CambiarColorFondo(slide, 204, 255, 204)

    # Crear diapositiva
    slide2 = CrearDiapositiva(prs, 2)

    # Cambiar fondo de la diapositiva
    CambiarColorFondo(slide2, 197, 25, 30)

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "presentacion.pptx")

    # Mensaje de presentación creada
    print("")
    print("La presentación se ha creado correctamente!")


if __name__ == '__main__':
    main()

