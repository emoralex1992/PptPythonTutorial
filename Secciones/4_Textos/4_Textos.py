from pptx import Presentation
import os
import shutil
import pandas as pd
import warnings
from pptx.util import Pt
import openpyxl
from pptx.dml.color import RGBColor
warnings.filterwarnings("ignore", category=UserWarning, module='openpyxl')

# ------------------- CONFIGURACIÓN DE USUARIO -------------------

# Ruta de entrada: Plantilla PowerPoint
INPUT_PPTX_TEMPLATE = './Inputs/Template_GCU.pptx'

# Ruta de entrada: Fichero Excel de configuración de la presentación
INPUT_EXCEL_CFG = './Inputs/Presentacion_Cfg.xlsx'

# Ruta de salida
OUTPUT_PATH = './Outputs'

# ------------------- /CONFIGURACIÓN DE USUARIO ------------------

# Eliminar y crear carpeta de nuevo
def RestaurarCarpeta(folder_path):
    # Verificar si la carpeta existe
    if os.path.exists(folder_path):
        # Eliminar la carpeta
        shutil.rmtree(folder_path)

    # Crear carpeta
    os.makedirs(folder_path)


# Crear diapositiva
def CrearDiapositiva(presentacion, layout_id):
    # Crear diapositiva
    slide_layout = presentacion.slide_layouts[layout_id]
    slide = presentacion.slides.add_slide(slide_layout)

    return slide


# Añadir texto a un placeholder de una diapositiva
def AñadirTextoPlaceholder(slide, placeholder_id, runs_list, nivel=0, paragraph_idx=0):
    # Crear el cuadro de texto (textframe)
    body_shape = slide.shapes.placeholders[placeholder_id]
    tf = body_shape.text_frame

    # Comprobar si ya hay texto en el text_frame
    if paragraph_idx == 0:
        # Utilizar el párrafo 0
        p = tf.paragraphs[0]
    else:
        # Añadir un nuevo párrafo
        p = tf.add_paragraph()

    # Asignar nivel
    p.level = nivel

    # Iterar sobre runs_list
    for run_dict in runs_list:
        AñadirRunAlParrafo(p, run_dict)


# Rutina para añadir runs al párrafo
def AñadirRunAlParrafo(paragraph, run_dict):
    run = paragraph.add_run()
    run.text = run_dict['texto']

    # Añadir negrita
    if "negrita" in run_dict:
        run.font.bold = run_dict["negrita"]

    # Añadir cursiva
    if "cursiva" in run_dict:
        run.font.italic = run_dict["cursiva"]

    # Añadir subrayado
    if "subrayado" in run_dict:
        run.font.underline = run_dict["subrayado"]

    # Añadir tamaño del texto
    if "tamaño_txt" in run_dict and run_dict['tamaño_txt'] != None:
        run.font.size = Pt(run_dict["tamaño_txt"])

    # Añadir color de fuente
    if "color_txt" in run_dict and run_dict['color_txt'] != None:
        run.font.color.rgb = RGBColor(*run_dict["color_txt"])

    # Añadir tipo de fuente
    if "fuente" in run_dict and run_dict['fuente'] != None:
        run.font.name = run_dict["fuente"]


# Imprimir todos los nombres de los diferentes diseños de diapositivas (layouts) de una plantilla PowerPoint
def ImprimirDiseñosPresentacion(prs):
    print("LISTA DE DISEÑOS DE DIAPOSITIVAS:")
    # Iterar sobre los layouts
    for layout in prs.slide_layouts:
        # Imprimir el nombre del layout
        print(layout.name)
    print("")


# Crear presentación con cada uno de los diseños y el nombre de los placeholders
def ExportarNombrePlaceholders():
    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Iterar sobre los layouts
    for layout_id, layout in enumerate(prs.slide_layouts):
        slide = CrearDiapositiva(prs, layout_id)
        # print(layout.name)

        # Iterar sobre los placeholders
        for ph in layout.placeholders:
            try:
                text = str(ph.placeholder_format.idx) + " - " + ph.name
                runs_list = [
                    {
                        'texto': text,
                    }
                ]
                AñadirTextoPlaceholder(slide, ph.placeholder_format.idx, runs_list, 0, 0)
                # print(f"   {ph.placeholder_format.idx} - {ph.name}")
            except KeyError:
                i = 0

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "Marcadores.pptx")

# Convertir color hexadecimal a RGB
def HexToRgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return [int(hex_color[i:i+2], 16) for i in (0, 2, 4)]

# Importar datos desde el fichero de configuración Excel
def ImportarDatosExcel(prs):
    # Inicializar lista slide_list
    slide_list = []

    # Leer fichero Excel
    df = pd.read_excel(INPUT_EXCEL_CFG, sheet_name='Presentación')

    # Leer fichero Excel utilizando openpyxl
    wb = openpyxl.load_workbook(INPUT_EXCEL_CFG)
    ws = wb['Presentación']

    # Iterar sobre las filas del DataFrame
    for index, row in df.iterrows():
        # Leer celdas de la fila
        nombre_layout = row['DISEÑO_DIAPOSITIVA']
        placeholder_id = row['NUM_MARCADOR']
        nivel = int(row['NIVEL_TEXTO'].replace("N", "")) if pd.notna(row['NIVEL_TEXTO']) else None
        texto_run = str(row['TEXTO']) if pd.notna(row['TEXTO']) else ""
        negrita = row['NEGRITA'] == "Sí"
        cursiva = row['CURSIVA'] == "Sí"
        subrayado = row['SUBRAYADO'] == "Sí"
        tamaño_texto = int(row['TAMAÑO_TXT']) if pd.notna(row['TAMAÑO_TXT']) else None
        fuente = row['FUENTE'] if pd.notna(row['FUENTE']) else None

        # Obtener el color de la celda correspondiente a COLOR_TXT
        cell = ws.cell(row=index + 2, column=df.columns.get_loc('COLOR_TXT') + 1)
        color_celda = cell.fill.fgColor
        r, g, b = None, None, None  # Inicializar valores de color

        if color_celda.type == "rgb" and isinstance(color_celda.rgb, str):  # Si es un color directo en formato hexadecimal
            color_hex = color_celda.rgb
            if len(color_hex) == 8 and color_hex != "00000000":  # Verificar si no es transparente o sin color asignado
                r, g, b = HexToRgb(color_hex[2:])  # Convertir a RGB

        elif color_celda.type == "theme":  # Si es un color de tema
            # Lanzar un error informando al usuario que debe cambiar el color de tema
            raise ValueError(f"El color en la fila {index + 2} columna 'COLOR_TXT' es un color de tema. "
                             f"Por favor, cambie el color a un valor explícito (no de tema) en el archivo Excel.")

        # Diapositiva encontrada
        if pd.notna(nombre_layout):
            # Encontrar el layout_id correspondiente al layout_name
            layout_id = None
            for i, layout in enumerate(prs.slide_layouts):
                if layout.name == nombre_layout:
                    layout_id = i
                    break

            # Error: layout_id no encontrado
            if layout_id is None:
                row_error = 2 + index
                raise ValueError(f"No se encontró un diseño de diapositiva con el nombre: {nombre_layout} - FILA {row_error} / COLUMNA \"DISEÑO_DIAPOSITIVA\"")

            print(f"{nombre_layout} - Idx: {layout_id}")

            # Crear diccionario para la diapositiva y añadirlo a slide_list
            slide_list.append({
                'layout_name': nombre_layout,
                'layout': layout_id,
                'placeholders': []
            })

        # Marcador (placeholder) encontrado
        elif pd.notna(placeholder_id):
            # Verificar si el placeholder_id es válido para el diseño de la diapositiva actual
            layout_id = slide_list[-1]['layout']
            valid_placeholder_ids = [ph.placeholder_format.idx for ph in prs.slide_layouts[layout_id].placeholders]

            # Error: el placeholder no se encuentra en el diseño de diapositiva
            if placeholder_id not in valid_placeholder_ids:
                row_error = 2 + index
                raise ValueError(f"No se encontró el marcaddor con ID {int(placeholder_id)} en el diseño de diapositiva \"{slide_list[-1]['layout_name']}\" - FILA {row_error} / COLUMNA \"NUM_MARCADOR\"")

            # Añadir el placeholder al último slide_dict en slide_list
            slide_list[-1]['placeholders'].append({
                'placeholder_id': placeholder_id,
                'elementos': []
            })

        # Nivel_Texto encontrado
        elif nivel != None:
            # Añadir el texto a los elementos del último placeholder
            slide_list[-1]['placeholders'][-1]['elementos'].append({
                'runs': [],
                'nivel': nivel
            })

        # Texto encontrado
        elif texto_run != "":
            # Añadir run a lista "runs"
            slide_list[-1]['placeholders'][-1]['elementos'][-1]['runs'].append(
                {
                    'texto': texto_run,
                    'negrita': negrita,
                    'cursiva': cursiva,
                    'subrayado': subrayado,
                    'tamaño_txt': tamaño_texto,
                    'color_txt': [r, g, b] if r != None else None,
                    'fuente': fuente
                })

    return slide_list


# Rutina principal
def main():
    # Eliminar y crear OUTPUT_PATH de nuevo
    RestaurarCarpeta(OUTPUT_PATH)

    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Imprimir los diseños de la presentación por consola
    ImprimirDiseñosPresentacion(prs)

    # Exportar el nombre de los placeholders en una nueva presentación
    ExportarNombrePlaceholders()

    # Importar datos desde fichero Excel
    slide_list = ImportarDatosExcel(prs)

    # Iteración sobre la lista de diccionarios SLIDE_LIST
    for slide_dict in slide_list:
        # Crear diapositiva
        layout_id = slide_dict['layout']
        slide = CrearDiapositiva(prs, layout_id)

        # Iterar sobre los placeholders
        for ph_dict in slide_dict['placeholders']:
            placeholder_id = ph_dict['placeholder_id']

            # Iterar sobre cada elemento del placeholder
            for elem_idx, elem_dict in enumerate(ph_dict['elementos']):
                # Añadir texto al placeholder
                nivel = elem_dict['nivel']
                runs_list = elem_dict['runs']
                AñadirTextoPlaceholder(slide, placeholder_id, runs_list, nivel, elem_idx)

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "PlantillasProyecto_Output.pptx")

    # Mensaje de presentación creada
    print("")
    print("La presentación se ha creado correctamente!")


if __name__ == '__main__':
    main()

