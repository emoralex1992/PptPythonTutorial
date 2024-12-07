from pptx import Presentation
import os
import shutil
import pandas as pd
import warnings
warnings.filterwarnings("ignore", category=UserWarning, module='openpyxl')

# ------------------- CONFIGURACIÓN DE USUARIO -------------------

# Ruta de entrada: Plantilla PowerPoint
INPUT_PPTX_TEMPLATE = './Inputs/PlantillaPyCoffee.pptx'

# Ruta de entrada: Fichero Excel de configuración de la presentación
INPUT_EXCEL_CFG = './Inputs/Presentacion_Cfg.xlsx'

# Ruta de salida
OUTPUT_PATH = '../../Outputs'

# ------------------- /CONFIGURACIÓN DE USUARIO ------------------

# Eliminar y crear carpeta de nuevo
def RestaurarCarpeta(folder_path):
    # Verificar si la carpeta existe
    if os.path.exists(folder_path):
        # Eliminar la carpeta
        shutil.rmtree(folder_path)

    # Crear carpeta
    os.makedirs(folder_path)


# Crear diapositiva
def CrearDiapositiva(presentacion, layout_id):
    # Crear diapositiva
    slide_layout = presentacion.slide_layouts[layout_id]
    slide = presentacion.slides.add_slide(slide_layout)

    return slide


# Añadir texto a un placeholder de una diapositiva
def AñadirTextoPlaceholder(slide, placeholder_id, texto, nivel=0):
    # Crear el cuadro de texto (textframe)
    body_shape = slide.shapes.placeholders[placeholder_id]
    tf = body_shape.text_frame

    # Comprobbar si ya hay texto en el text_frame
    if tf.text:
        # Añadir un nuevo párrafo
        p = tf.add_paragraph()
        p.text = texto
        p.level = nivel
    else:
        # Reemplazar el texto existente
        tf.text = texto
        tf.level = nivel


# Imprimir todos los nombres de los diferentes diseños de diapositivas (layouts) de una plantilla PowerPoint
def ImprimirDiseñosPresentacion(prs):
    print("LISTA DE DISEÑOS DE DIAPOSITIVAS:")
    # Iterar sobre los layouts
    for layout in prs.slide_layouts:
        # Imprimir el nombre del layout
        print(layout.name)
    print("")


# Crear presentación con cada uno de los diseños y el nombre de los placeholders
def ExportarNombrePlaceholders():
    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Iterar sobre los layouts
    for layout_id, layout in enumerate(prs.slide_layouts):
        slide = CrearDiapositiva(prs, layout_id)
        # print(layout.name)

        # Iterar sobre los placeholders
        for ph in layout.placeholders:
            try:
                text = str(ph.placeholder_format.idx) + " - " + ph.name
                AñadirTextoPlaceholder(slide, ph.placeholder_format.idx, text, 0)
                # print(f"   {ph.placeholder_format.idx} - {ph.name}")
            except KeyError:
                i = 0

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "Marcadores.pptx")


# Importar datos desde el fichero de configuración Excel
def ImportarDatosExcel(prs):
    # Inicializar lista slide_list
    slide_list = []

    # Leer fichero Excel
    df = pd.read_excel(INPUT_EXCEL_CFG, sheet_name='Presentación')

    # Iterar sobre las filas del DataFrame
    for index, row in df.iterrows():
        # Leer celdas de la fila
        nombre_layout = row['DISEÑO_DIAPOSITIVA']
        placeholder_id = row['NUM_MARCADOR']
        nivel = int(row['NIVEL_TEXTO']) if pd.notna(row['NIVEL_TEXTO']) else 0
        texto = str(row['TEXTO']) if pd.notna(row['TEXTO']) else ""

        # Diapositiva encontrada
        if pd.notna(nombre_layout):
            # Encontrar el layout_id correspondiente al layout_name
            layout_id = None
            for i, layout in enumerate(prs.slide_layouts):
                if layout.name == nombre_layout:
                    layout_id = i
                    break

            # Error: layout_id no encontrado
            if layout_id is None:
                row_error = 2 + index
                raise ValueError(f"No se encontró un diseño de diapositiva con el nombre: {nombre_layout} - FILA {row_error} / COLUMNA \"DISEÑO_DIAPOSITIVA\"")

            print(f"{nombre_layout} - Idx: {layout_id}")

            # Crear diccionario para la diapositiva y añadirlo a slide_list
            slide_list.append({
                'layout_name': nombre_layout,
                'layout': layout_id,
                'placeholders': []
            })
        # Marcador (placeholder) encontrado
        elif pd.notna(placeholder_id):
            # Verificar si el placeholder_id es válido para el diseño de la diapositiva actual
            layout_id = slide_list[-1]['layout']
            valid_placeholder_ids = [ph.placeholder_format.idx for ph in prs.slide_layouts[layout_id].placeholders]

            # Error: el placeholder no se encuentra en el diseño de diapositiva
            if placeholder_id not in valid_placeholder_ids:
                row_error = 2 + index
                raise ValueError(f"No se encontró el marcaddor con ID {int(placeholder_id)} en el diseño de diapositiva \"{slide_list[-1]['layout_name']}\" - FILA {row_error} / COLUMNA \"NUM_MARCADOR\"")

            # Añadir el placeholder al último slide_dict en slide_list
            slide_list[-1]['placeholders'].append({
                'placeholder_id': placeholder_id,
                'elementos': []
            })
        # Texto encontrado
        elif pd.notna(nivel):
            # Añadir el texto a los elementos del último placeholder
            slide_list[-1]['placeholders'][-1]['elementos'].append({
                'texto': texto,
                'nivel': nivel
            })

    return slide_list


# Rutina principal
def main():
    # Eliminar y crear OUTPUT_PATH de nuevo
    RestaurarCarpeta(OUTPUT_PATH)

    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Imprimir los diseños de la presentación por consola
    ImprimirDiseñosPresentacion(prs)

    # Exportar el nombre de los placeholders en una nueva presentación
    ExportarNombrePlaceholders()

    # Importar datos desde fichero Excel
    slide_list = ImportarDatosExcel(prs)

    # Iteración sobre la lista de diccionarios SLIDE_LIST
    for slide_dict in slide_list:
        # Crear diapositiva
        layout_id = slide_dict['layout']
        slide = CrearDiapositiva(prs, layout_id)

        # # Modificar el título de la diapositiva (placeholder 0)
        # texto = slide_dict['texto_titulo']
        # AñadirTextoPlaceholder(slide, 0, texto)

        # Iterar sobre los placeholders
        for ph_dict in slide_dict['placeholders']:
            placeholder_id = ph_dict['placeholder_id']

            # Iterar sobre cada elemento del placeholder
            for elem_dict in ph_dict['elementos']:
                # Añadir texto al placeholder
                nivel = elem_dict['nivel']
                texto = elem_dict['texto']
                AñadirTextoPlaceholder(slide, placeholder_id, texto, nivel)

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "PlantillasProyecto_Output.pptx")

    # Mensaje de presentación creada
    print("")
    print("La presentación se ha creado correctamente!")


if __name__ == '__main__':
    main()

