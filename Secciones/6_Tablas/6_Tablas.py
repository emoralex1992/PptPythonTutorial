from pptx import Presentation
import os
import shutil
import pandas as pd
import warnings
from pptx.util import Pt, Cm
import openpyxl
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from datetime import datetime
import re
import locale

# Configurar el locale a español (es_ES para España)
locale.setlocale(locale.LC_TIME, 'es_ES.UTF-8')

# Ignorar warnings
warnings.filterwarnings("ignore", category=UserWarning, module='openpyxl')

# ------------------- CONFIGURACIÓN DE USUARIO -------------------

# Ruta de entrada: Plantilla PowerPoint
INPUT_PPTX_TEMPLATE = './Inputs/Template_GCU.pptx'

# Ruta de entrada: Fichero Excel de configuración de la presentación
INPUT_EXCEL_CFG = './Inputs/Presentacion_Cfg.xlsx'

# Ruta de salida
OUTPUT_PATH = './Outputs'

# ------------------- /CONFIGURACIÓN DE USUARIO ------------------

# -------------------------- CONSTANTES --------------------------

TABLE_STYLE_ID_DICT = {
    "SinEstiloSinContorno":         '{2D5ABB26-0587-4C30-8999-92F81FD0307C}',
    "EstiloTematico1Acento1":       '{3C2FFA5D-87B4-456A-9821-1D50468CF0F}',
    "EstiloTematico1Acento2":       '{284E427A-3D55-4303-BF80-6455036E1DE7}',
    "EstiloTematico1Acento3":       '{69C7853C-536D-4A76-A0AE-DD22124D55A5}',
    "EstiloTematico1Acento4":       '{775DCB02-9BB8-47FD-8907-85C794F793BA}',
    "EstiloTematico1Acento5":       '{35758FB7-9AC5-4552-8A53-C91805E547FA}',
    "EstiloTematico1Acento6":       '{08FB837D-C827-4EFA-A057-4D05807E0F7C}',
    "SinEstiloContornoTabla":       '{5940675A-B579-460E-94D1-54222C63F5DA}',
    "EstiloTematico2Acento1":       '{D113A9D2-9D6B-4929-AA2D-F23B5EE8CBE7}',
    "EstiloTematico2Acento2":       '{18603FDC-E32A-4AB5-989C-0864C3EAD2B8}',
    "EstiloTematico2Acento3":       '{306799F8-075E-4A3A-A7F6-7FBC6576F1A4}',
    "EstiloTematico2Acento4":       '{E269D01E-BC32-4049-B463-5C60D7B0CCD2}',
    "EstiloTematico2Acento5":       '{327F97BB-C833-4FB7-BDE5-3F7075034690}',
    "EstiloTematico2Acento6":       '{638B1855-1B75-4FBE-930C-398BA8C253C6}',
    "EstiloClaro1":                 '{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}',
    "EstiloClaro1Acento1":          '{3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5}',
    "EstiloClaro1Acento2":          '{0E3FDE45-AF77-4B5C-9715-49D594BDF05E}',
    "EstiloClaro1Acento3":          '{C083E6E3-FA7D-4D7B-A595-EF9225AFEA82}',
    "EstiloClaro1Acento4":          '{D27102A9-8310-4765-A935-A1911B00CA55}',
    "EstiloClaro1Acento5":          '{5FD0F851-EC5A-4D38-B0AD-8093EC10F338}',
    "EstiloClaro1Acento6":          '{68D230F3-CF80-4859-8CE7-A43EE81993B5}',
    "EstiloClaro2":                 '{7E9639D4-E3E2-4D34-9284-5A2195B3D0D7}',
    "EstiloClaro2Acento1":          '{69012ECD-51FC-41F1-AA8D-1B2483CD663E}',
    "EstiloClaro2Acento2":          '{72833802-FEF1-4C79-8D5D-14CF1EAF98D9}',
    "EstiloClaro2Acento3":          '{F2DE63D5-997A-4646-A377-4702673A728D}',
    "EstiloClaro2Acento4":          '{17292A2E-F333-43FB-9621-5CBBE7FDCDCB}',
    "EstiloClaro2Acento5":          '{5A111915-BE36-4E01-A7E5-04B1672EAD32}',
    "EstiloClaro2Acento6":          '{912C8C85-51F0-491E-9774-3900AFEF0FD7}',
    "EstiloClaro3":                 '{616DA210-FB5B-4158-B5E0-FEB733F419BA}',
    "EstiloClaro3Acento1":          '{BC89EF96-8CEA-46FF-86C4-4CE0E7609802}',
    "EstiloClaro3Acento2":          '{5DA37D80-6434-44D0-A028-1B22A696006F}',
    "EstiloClaro3Acento3":          '{8799B23B-EC83-4686-B30A-512413B5E67A}',
    "EstiloClaro3Acento4":          '{ED083AE6-46FA-4A59-8FB0-9F97EB10719F}',
    "EstiloClaro3Acento5":          '{BDBED569-4797-4DF1-A0F4-6AAB3CD982D8}',
    "EstiloClaro3Acento6":          '{E8B1032C-EA38-4F05-BA0D-38AFFFC7BED3}',
    "EstiloMedio1":                 '{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}',
    "EstiloMedio1Acento1":          '{B301B821-A1FF-4177-AEE7-76D212191A09}',
    "EstiloMedio1Acento2":          '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}',
    "EstiloMedio1Acento3":          '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}',
    "EstiloMedio1Acento4":          '{1E171933-4619-4E11-9A3F-F7608DF75F80}',
    "EstiloMedio1Acento5":          '{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}',
    "EstiloMedio1Acento6":          '{10A1B5D5-9B99-4C35-A422-299274C87663}',
    "EstiloMedio2":                 '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}',
    "EstiloMedio2Acento1":          '{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}',
    "EstiloMedio2Acento2":          '{21E4AEA4-8DFA-4A89-87EB-49C32662AFE0}',
    "EstiloMedio2Acento3":          '{F5AB1C69-6EDB-4FF4-983F-18BD219EF322}',
    "EstiloMedio2Acento4":          '{00A15C55-8517-42AA-B614-E9B94910E393}',
    "EstiloMedio2Acento5":          '{7DF18680-E054-41AD-8BC1-D1AEF772440D}',
    "EstiloMedio2Acento6":          '{93296810-A885-4BE3-A3E7-6D5BEEA58F35}',
    "EstiloMedio3":                 '{8EC20E35-A176-4012-BC5E-935CFFF8708E}',
    "EstiloMedio3Acento1":          '{6E25E649-3F16-4E02-A733-19D2CDBF48F0}',
    "EstiloMedio3Acento2":          '{85BE263C-DBD7-4A20-BB59-AAB30ACAA65A}',
    "EstiloMedio3Acento3":          '{EB344D84-9AFB-497E-A393-DC336BA19D2E}',
    "EstiloMedio3Acento4":          '{EB9631B5-78F2-41C9-869B-9F39066F8104}',
    "EstiloMedio3Acento5":          '{74C1A8A3-306A-4EB7-A6B1-4F7E0EB9C5D6}',
    "EstiloMedio3Acento6":          '{2A488322-F2BA-4B5B-9748-0D474271808F}',
    "EstiloMedio4":                 '{D7AC3CCA-C797-4891-BE02-D94E43425B78}',
    "EstiloMedio4Acento1":          '{69CF1AB2-1976-4502-BF36-3FF5EA218861}',
    "EstiloMedio4Acento2":          '{8A107856-5554-42FB-B03E-39F5DBC370BA}',
    "EstiloMedio4Acento3":          '{0505E3EF-67EA-436B-97B2-0124C06EBD24}',
    "EstiloMedio4Acento4":          '{C4B1156A-380E-4F78-BDF5-A606A8083BF9}',
    "EstiloMedio4Acento5":          '{22838BEF-8BB2-4498-84A7-C5851F593DF1}',
    "EstiloMedio4Acento6":          '{16D9F66E-5EB9-4882-86FB-DCBF35E3C3E4}',
    "EstiloOscuro1":                '{E8034E78-7F5D-4C2E-B375-FC64B27BC917}',
    "EstiloOscuro1Acento1":         '{125E5076-3810-47DD-B79F-674D7AD40C01}',
    "EstiloOscuro1Acento2":         '{37CE84F3-28C3-443E-9E96-99CF82512B78}',
    "EstiloOscuro1Acento3":         '{D03447BB-5D67-496B-8E87-E561075AD55C}',
    "EstiloOscuro1Acento4":         '{E929F9F4-4A8F-4326-A1B4-22849713DDAB}',
    "EstiloOscuro1Acento5":         '{8FD4443E-F989-4FC4-A0C8-D5A2AF1F390B}',
    "EstiloOscuro1Acento6":         '{AF606853-7671-496A-8E4F-DF71F8EC918B}',
    "EstiloOscuro2":                '{5202B0CA-FC54-4496-8BCA-5EF66A818D29}',
    "EstiloOscuro2Acento1Acento2":  '{0660B408-B3CF-4A94-85FC-2B1E0A45F4A2}',
    "EstiloOscuro2Acento3Acento4":  '{91EBBBCC-DAD2-459C-BE2E-F6DE35CF9A28}',
    "EstiloOscuro2Acento5Acento6":  '{46F890A9-2807-4EBB-B81D-B2AA78EC7F39}'
}

# -------------------------- /CONSTANTES -------------------------

# Eliminar y crear carpeta de nuevo
def RestaurarCarpeta(folder_path):
    # Verificar si la carpeta existe
    if os.path.exists(folder_path):
        # Eliminar la carpeta
        shutil.rmtree(folder_path)

    # Crear carpeta
    os.makedirs(folder_path)


# Crear diapositiva
def CrearDiapositiva(presentacion, layout_id):
    # Crear diapositiva
    slide_layout = presentacion.slide_layouts[layout_id]
    slide = presentacion.slides.add_slide(slide_layout)

    return slide


# Añadir texto a un placeholder de una diapositiva
def AñadirTextoPlaceholder(slide, placeholder_id, runs_list, nivel=0, paragraph_idx=0):
    # Crear el cuadro de texto (textframe)
    body_shape = slide.shapes.placeholders[placeholder_id]
    tf = body_shape.text_frame

    # Comprobar si ya hay texto en el text_frame
    if paragraph_idx == 0:
        # Utilizar el párrafo 0
        p = tf.paragraphs[0]
    else:
        # Añadir un nuevo párrafo
        p = tf.add_paragraph()

    # Asignar nivel
    p.level = nivel

    # Iterar sobre runs_list
    for run_dict in runs_list:
        AñadirRunAlParrafo(p, run_dict)


# Rutina para añadir runs al párrafo
def AñadirRunAlParrafo(paragraph, run_dict):
    run = paragraph.add_run()
    run.text = run_dict['texto']

    # Añadir negrita
    if "negrita" in run_dict:
        run.font.bold = run_dict["negrita"]

    # Añadir cursiva
    if "cursiva" in run_dict:
        run.font.italic = run_dict["cursiva"]

    # Añadir subrayado
    if "subrayado" in run_dict:
        run.font.underline = run_dict["subrayado"]

    # Añadir tamaño del texto
    if "tamaño_txt" in run_dict and run_dict['tamaño_txt'] != None:
        run.font.size = Pt(run_dict["tamaño_txt"])

    # Añadir color de fuente
    if "color_txt" in run_dict and run_dict['color_txt'] != None:
        run.font.color.rgb = RGBColor(*run_dict["color_txt"])

    # Añadir tipo de fuente
    if "fuente" in run_dict and run_dict['fuente'] != None:
        run.font.name = run_dict["fuente"]


# Rutina para añadir una imagen a un placeholder
def AñadirImagenPlaceholder(slide, img_path, placeholder, img_param_dict):
    # Insertar imagen en la diapositiva
    img = slide.shapes.add_picture(img_path, placeholder.left, placeholder.top)

    # Calcular la relación de aspecto del marcador de posición
    placeholder_aspect_ratio = placeholder.width / placeholder.height

    # Calcular la relación de aspecto de la imagen
    img_aspect_ratio = img.width / img.height

    # Recortar imagen de forma automática para ajustarla al placeholder
    if img_param_dict['img_modo'] == "RECORTAR PARA ENCAJAR":
        # Ajustar el ancho y alto de la imagen al ancho y alto del placeholder
        img.height = placeholder.height
        img.width = placeholder.width

        # La imagen es más ancha que el placeholder
        if img_aspect_ratio > placeholder_aspect_ratio:
            # Calcular la proporción de anchura
            image_prop_width = int(img.height * img_aspect_ratio)

            # Recortar la imagen para ajustarla a la anchura del placeholder basándonos en la anchura proporcional
            crop_value = ((image_prop_width - img.width) / image_prop_width) / 2
            img.crop_left = crop_value
            img.crop_right = crop_value

        # La imagen es más alta que el placeholder
        else:
            # Calcular la proporción de altura
            image_prop_height = int(img.width / img_aspect_ratio)

            # Recortar la imagen para ajustarla a la altura del placeholder basándonos en la altura proporcional
            crop_value = ((image_prop_height - img.height) / image_prop_height) / 2
            img.crop_top = crop_value
            img.crop_bottom = crop_value

        # Aplicar offset al recorte
        img.crop_left += img_param_dict['rec_offset_der_porc'] - img_param_dict['rec_offset_izq_porc']
        img.crop_right += img_param_dict['rec_offset_izq_porc'] - img_param_dict['rec_offset_der_porc']
        img.crop_top += img_param_dict['rec_offset_inf_porc'] - img_param_dict['rec_offset_sup_porc']
        img.crop_bottom += img_param_dict['rec_offset_sup_porc'] - img_param_dict['rec_offset_inf_porc']

    # Imagen con ancho y alto automáticos --> Ajustar al placeholder
    elif img_param_dict['img_modo'] == "ANCHO/ALTO AUTOMATICO":
        # La imagen es más ancha que el placeholder
        if img_aspect_ratio > placeholder_aspect_ratio:
            scale_factor = placeholder.width / img.width
        # La imagen es más alta que el placeholder
        else:
            scale_factor = placeholder.height / img.height

        # Aplicar el factor de escala a la imagen
        img.width = int(img.width * scale_factor)
        img.height = int(img.height * scale_factor)

    # Imagen con ancho y alto personalizados
    elif img_param_dict['img_modo'] == "TAMAÑO PERSONALIZADO":
        # Ancho y Alto definidos
        if (img_param_dict['ancho_img_cm'] != None) and (img_param_dict['alto_img_cm'] != None):
            img.width = Cm(img_param_dict['ancho_img_cm'])
            img.height = Cm(img_param_dict['alto_img_cm'])

        # Únicamente el Ancho de la imagen ha sido definido
        elif (img_param_dict['ancho_img_cm'] != None) and (img_param_dict['alto_img_cm'] == None):
            img.height = Cm(img_param_dict['ancho_img_cm'] * (img.height / img.width))
            img.width = Cm(img_param_dict['ancho_img_cm'])

        # Únicamente el Alto de la imagen ha sido definido
        elif (img_param_dict['ancho_img_cm'] == None) and (img_param_dict['alto_img_cm'] != None):
            img.width = Cm(img_param_dict['alto_img_cm'] * (img.width / img.height))
            img.heigh = Cm(img_param_dict['alto_img_cm'])

        # Error: Ni el Alto ni el Ancho de la imagen han sido definidos
        else:
            raise ValueError(f"El modo de imagen TAMAÑO PERSONALIZADO debe tener uno o todos los campos \"ALTO_IMG_CM\" y/o \"ANCHO_IMG_CM\" definidos "
                             f"para la imagen: {img_path}")

    # Alinear la imagen verticalmente dentro del placeholder ("IZQUIERDA" por defecto)
    if img_param_dict['img_alinear_v'] == "DERECHA":
        width_diff = placeholder.width - img.width
        img.left = placeholder.left + width_diff
    elif img_param_dict['img_alinear_v'] == "CENTRO":
        width_diff = placeholder.width - img.width
        img.left = placeholder.left + int(width_diff / 2)

    # Alinear la imagen horizontalmente dentro del placeholder ("ARRIBA" por defecto)
    if img_param_dict['img_alinear_h'] == "ABAJO":
        height_diff = placeholder.height - img.height
        img.top = placeholder.top + height_diff
    elif img_param_dict['img_alinear_h'] == "CENTRO":
        height_diff = placeholder.height - img.height
        img.top = placeholder.top + int(height_diff / 2)


# Rutina para extraer y aplicar formato de texto a una celda
def AñadirTextoFormatoCelda(cell, raw_text):
    # Expresión regular para buscar los formatos avanzados
    patron = re.compile(r'{(.*?)}')

    # Inicializar valores de formato
    texto = str(raw_text)
    formato = {}

    # Buscar y estraer formatos avanzados si existen en el valor de la celda
    match = patron.search(texto)
    if match:
        # Extraer el contenido dentro de los corchetes
        formato_str = match.group(1)

        # Separar las propiedades del formato
        formato = {item.split('=')[0].strip(): item.split('=')[1].strip() for item in formato_str.split(', ')}

        # Quitar los corchetes del valor de la celda para obtener el texto limpio
        texto = patron.sub('', texto).strip()

    # Añadir el texto limpio a la celda
    cell.text = texto if pd.notna(raw_text) else ''

    # Aplicar el formato a cada párrafo de la celda
    for paragraph in cell.text_frame.paragraphs:
        # Aplicar negrita si está especificado en el formato
        if 'negrita' in formato and formato['negrita'] == 'Sí':
            paragraph.font.bold = True

        # Aplicar cursiva si está especificado en el formato
        if 'cursiva' in formato and formato['cursiva'] == 'Sí':
            paragraph.font.italic = True

        # Aplicar subrayado si está especificado en el formato
        if 'subrayado' in formato and formato['subrayado'] == 'Sí':
            paragraph.font.underline = True

        # Aplicar tamaño si está especificado en el formato
        if 'tamaño' in formato:
            try:
                paragraph.font.size = Pt(float(formato['tamaño']))
            except ValueError:
                print(f"Error: Tamaño de texto no válido en {formato['tamaño']}")

        # Aplicar color si está especificado en el formato
        if 'color' in formato:
            color_rgb = HexToRgb(formato['color'])
            paragraph.font.color.rgb = RGBColor(*color_rgb)


# Rutina para añadir una tabla a un placeholder
def AñadirTablaPlaceholder(slide, placeholder_idx, elem_params):
    # Obtener placeholder
    placeholder = slide.placeholders[placeholder_idx]

    # Obtener parámetros del diccionario elem_params
    tabla = elem_params['tabla']
    estilo_tabla = elem_params['estilo_tabla']
    tamaño_texto = elem_params['tamaño_txt']
    fuente = elem_params['fuente']

    # Calcular número de filas y columnas de la tabla
    nrows = len(tabla)
    ncols = len(tabla[0]) if nrows > 0 else 0

    # Crear la tabla en el placeholder
    table_shape = slide.shapes.add_table(nrows, ncols, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
    table = table_shape.table

    # Añadir estilo a la tabla
    tbl = table_shape._element.graphic.graphicData.tbl
    style_id = TABLE_STYLE_ID_DICT[estilo_tabla]
    tbl[0][-1].text = style_id

    # Añadir valores a la tabla
    for row_idx, row in enumerate(tabla):
        for col_idx, cell_value in enumerate(row):
            # Obtener celda
            cell = table.cell(row_idx, col_idx)

            # Añadir texto a la celda
            AñadirTextoFormatoCelda(cell, cell_value)

            # Aplicar tamaño de texto y fuente si están definidos
            for paragraph in cell.text_frame.paragraphs:
                if tamaño_texto and not(paragraph.font.size):
                    paragraph.font.size = Pt(tamaño_texto)
                if fuente and not(paragraph.font.name):
                    paragraph.font.name = fuente

                # Alinear parágrafo horizontalmente al centro
                paragraph.alignment = PP_ALIGN.CENTER

            # Alinear celda verticalmente al centro
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# Imprimir todos los nombres de los diferentes diseños de diapositivas (layouts) de una plantilla PowerPoint
def ImprimirDiseñosPresentacion(prs):
    print("LISTA DE DISEÑOS DE DIAPOSITIVAS:")
    # Iterar sobre los layouts
    for layout in prs.slide_layouts:
        # Imprimir el nombre del layout
        print(layout.name)
    print("")


# Crear presentación con cada uno de los diseños y el nombre de los placeholders
def ExportarNombrePlaceholders():
    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Iterar sobre los layouts
    for layout_id, layout in enumerate(prs.slide_layouts):
        slide = CrearDiapositiva(prs, layout_id)
        # print(layout.name)

        # Iterar sobre los placeholders
        for ph in layout.placeholders:
            try:
                text = str(ph.placeholder_format.idx) + " - " + ph.name
                runs_list = [
                    {
                        'texto': text,
                    }
                ]
                AñadirTextoPlaceholder(slide, ph.placeholder_format.idx, runs_list, 0, 0)
                # print(f"   {ph.placeholder_format.idx} - {ph.name}")
            except KeyError:
                i = 0

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "Marcadores.pptx")

# Convertir color hexadecimal a RGB
def HexToRgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return [int(hex_color[i:i+2], 16) for i in (0, 2, 4)]

# Rutina para extraer los formatos avanzados de los títulos y devolverlos en forma de diccionario
def ExtraerFormatosAvanzados(df):
    formatos = {}
    nuevos_titulos = {}

    # Expresión regular para buscar los formatos avanzados
    patron = re.compile('{(.*?)}')

    # Iterar sobre las columnas del DataFrame
    for titulo in df.columns:
        # Buscar un formato dentro de los corchetes { }
        match = patron.search(titulo)

        # Formato encontrado
        if match:
            # Extraer el contenido dentro de los corchetes
            formato_str = match.group(1)

            opciones = {}
            # Dividir las opciones si tienen clave-valor (fecha por defecto en caso de no tener clave-valor)
            for item in formato_str.split(', '):
                if '=' in item:
                    clave, valor = item.split('=', 1)
                    opciones[clave] = valor
                else:
                    # Mantener casos sin clave-valor, como fomato de fecha con texto
                    opciones['fecha'] = item

            # Almacenar el formato en el diccionario de formatos
            formatos[titulo] = opciones

            # Quitar la parte entre corchetes y el espacio anterior si lo hay
            nuevo_titulo = patron.sub('', titulo).strip()

            # Cambiar el título original por el modificado (sin corchetes)
            nuevos_titulos[titulo] = nuevo_titulo

        # No se ha encontrado ningún formato
        else:
            formatos[titulo] = None
            nuevos_titulos[titulo] = titulo     # Título sin cambios

    # Renombrar las columnas del DataFrame
    df.rename(columns=nuevos_titulos, inplace=True)

    # Devolver los formatos y el DataFrame actualizado
    return formatos, df


# Rutina para convertir notaciones de fecha simplificadas a códigos de strftime
def TraducirFormatoFecha(formato_simplificado):
    formatos_fecha = {
        'dd/mm/yyyy': '%d/%m/%Y',
        'mm/dd/yyyy': '%m/%d/%Y',
        'dd-mmm-yyyy': '%d-%b-%Y',
        'yyyy-mm-dd': '%Y-%m-%d',
        'd de mmm de yyyy': '%d de %B de %Y',
        'dd/mmm/yyyy hh:mm': '%d/%b/%Y %H:%M',
        'dd-mmm-yyyy hh:mm': '%d-%b-%Y %H:%M',
    }

    return formatos_fecha.get(formato_simplificado, '%d/%m/%Y') # Predeterminado a 'dd/mm/yyyy'

# Rutina para aplicar un formato avanzado (fecha, decimales, símbolo y posición de las unidades)
def AplicarFormatoAvanzado(valor, formato):
    if formato and pd.notna(valor):
        # Comprobación para formato de fecha
        if 'fecha' in formato:
            # Traducir formato de fecha simplificado a strftime
            formato_fecha = TraducirFormatoFecha(formato['fecha'])
            if isinstance(valor, (pd.Timestamp, datetime)):
                return valor.strftime(formato_fecha)

        # Intentar convertir el valor a float para números
        try:
            # Expresión regular para buscar los formatos avanzados de la celda
            patron = re.compile(r'{(.*?)}')
            # Buscar y extraer formatos avanzados si existen en el valor de la celda
            match = patron.search(str(valor))
            # Inicializamos las variables
            valor_sin_formato = valor
            texto_formato = ''
            if match:
                # Quitar los corchetes del valor de la celda para obtener el texto limpio
                valor_sin_formato = patron.sub('', valor).strip()
                texto_formato = valor.replace(valor_sin_formato, '')

                # Remplazar comas por puntos para que identifique el número como "float"
                valor_sin_formato = valor_sin_formato.replace(',', '.')

            valor_float = float(valor_sin_formato)

            # Aplicar el número de decimales si se ha especificado
            decimales = formato.get('decimales', 0)
            valor = f"{valor_float:.{decimales}f}" # Aplicar el número de decimales

            # Determinal ros separadores de miles y decimales
            sep_decimal = formato.get('sep_decimal', '.') # Punto por defecto
            sep_miles = ',' if sep_decimal == '.' else '.'

            # Formatear el número con separadores de miles y decimales
            if sep_decimal == ',':
                valor_formateado = f"{valor_float:,.{decimales}f}".replace(',', 'X').replace('.', ',').replace('X', sep_miles)
            else: # Separador decimal = '.'
                valor_formateado = f"{valor_float:,.{decimales}f}".replace(',', sep_miles)

            #Formatear el valor según la posición del símbolo
            if formato['posicion'] == 'f': # Símbolo al final
                return f"{valor_formateado}{formato['simbolo']}{texto_formato}"
            elif formato['posicion'] == 'i': # Símbolo al inicio
                return f"{formato['simbolo']}{valor_formateado}{texto_formato}"

            return valor_formateado

        except ValueError:
            # Si no se puede convertir a float, devolver el valor original
            pass

    # Si no hay formato o no se puede convertir, devolver el valor original
    return valor

# Importar datos desde el fichero de configuración Excel
def ImportarDatosExcel(prs):
    # Inicializar lista slide_list
    slide_list = []

    # Leer fichero Excel
    df = pd.read_excel(INPUT_EXCEL_CFG, sheet_name='Presentación')

    # Leer fichero Excel utilizando openpyxl
    wb = openpyxl.load_workbook(INPUT_EXCEL_CFG)
    ws = wb['Presentación']

    # Iterar sobre las filas del DataFrame
    for index, row in df.iterrows():
        # Leer celdas de la fila
        nombre_layout = row['DISEÑO_DIAPOSITIVA']
        placeholder_id = row['NUM_MARCADOR']
        nivel = int(row['NIVEL_TEXTO'].replace("N", "")) if pd.notna(row['NIVEL_TEXTO']) else None
        texto_run = str(row['TEXTO']) if pd.notna(row['TEXTO']) else ""
        negrita = row['NEGRITA'] == "Sí"
        cursiva = row['CURSIVA'] == "Sí"
        subrayado = row['SUBRAYADO'] == "Sí"
        tamaño_texto = int(row['TAMAÑO_TXT']) if pd.notna(row['TAMAÑO_TXT']) else None
        fuente = row['FUENTE'] if pd.notna(row['FUENTE']) else None
        img_path = row['IMG_PATH'] if pd.notna(row['IMG_PATH']) else None
        img_modo = row['IMG_MODO'] if pd.notna(row['IMG_MODO']) else None
        rec_offset_izq_porc = float(row['REC_OFFSET_IZQ_PORC']) if pd.notna(row['REC_OFFSET_IZQ_PORC']) else 0.0
        rec_offset_der_porc = float(row['REC_OFFSET_DER_PORC']) if pd.notna(row['REC_OFFSET_DER_PORC']) else 0.0
        rec_offset_sup_porc = float(row['REC_OFFSET_SUP_PORC']) if pd.notna(row['REC_OFFSET_SUP_PORC']) else 0.0
        rec_offset_inf_porc = float(row['REC_OFFSET_INF_PORC']) if pd.notna(row['REC_OFFSET_INF_PORC']) else 0.0
        ancho_img_cm = float(row['ANCHO_IMG_CM']) if pd.notna(row['ANCHO_IMG_CM']) else None
        alto_img_cm = float(row['ALTO_IMG_CM']) if pd.notna(row['ALTO_IMG_CM']) else None
        img_alinear_v = row['IMG_ALINEAR_V'] if pd.notna(row['IMG_ALINEAR_V']) else None
        img_alinear_h = row['IMG_ALINEAR_H'] if pd.notna(row['IMG_ALINEAR_H']) else None
        tabla = row['TABLA'] if pd.notna(row['TABLA']) else None
        estilo_tabla = row['ESTILO_TABLA'] if pd.notna(row['ESTILO_TABLA']) else None

        # Obtener el color de la celda correspondiente a COLOR_TXT
        cell = ws.cell(row=index + 2, column=df.columns.get_loc('COLOR_TXT') + 1)
        color_celda = cell.fill.fgColor
        r, g, b = None, None, None  # Inicializar valores de color

        if color_celda.type == "rgb" and isinstance(color_celda.rgb, str):  # Si es un color directo en formato hexadecimal
            color_hex = color_celda.rgb
            if len(color_hex) == 8 and color_hex != "00000000":  # Verificar si no es transparente o sin color asignado
                r, g, b = HexToRgb(color_hex[2:])  # Convertir a RGB

        elif color_celda.type == "theme":  # Si es un color de tema
            # Lanzar un error informando al usuario que debe cambiar el color de tema
            raise ValueError(f"El color en la fila {index + 2} columna 'COLOR_TXT' es un color de tema. "
                             f"Por favor, cambie el color a un valor explícito (no de tema) en el archivo Excel.")

        # Diapositiva encontrada
        if pd.notna(nombre_layout):
            # Encontrar el layout_id correspondiente al layout_name
            layout_id = None
            for i, layout in enumerate(prs.slide_layouts):
                if layout.name == nombre_layout:
                    layout_id = i
                    break

            # Error: layout_id no encontrado
            if layout_id is None:
                row_error = 2 + index
                raise ValueError(f"No se encontró un diseño de diapositiva con el nombre: {nombre_layout} - FILA {row_error} / COLUMNA \"DISEÑO_DIAPOSITIVA\"")

            print(f"{nombre_layout} - Idx: {layout_id}")

            # Crear diccionario para la diapositiva y añadirlo a slide_list
            slide_list.append({
                'layout_name': nombre_layout,
                'layout': layout_id,
                'placeholders': []
            })

        # Marcador (placeholder) encontrado
        elif pd.notna(placeholder_id):
            # Verificar si el placeholder_id es válido para el diseño de la diapositiva actual
            layout_id = slide_list[-1]['layout']
            valid_placeholder_ids = [ph.placeholder_format.idx for ph in prs.slide_layouts[layout_id].placeholders]

            # Error: el placeholder no se encuentra en el diseño de diapositiva
            if placeholder_id not in valid_placeholder_ids:
                row_error = 2 + index
                raise ValueError(f"No se encontró el marcaddor con ID {int(placeholder_id)} en el diseño de diapositiva \"{slide_list[-1]['layout_name']}\" - FILA {row_error} / COLUMNA \"NUM_MARCADOR\"")

            # Añadir el placeholder al último slide_dict en slide_list
            slide_list[-1]['placeholders'].append({
                'placeholder_id': placeholder_id,
                'elementos': []
            })

        # Nivel_Texto encontrado
        elif nivel != None:
            # Añadir el texto a los elementos del último placeholder
            slide_list[-1]['placeholders'][-1]['elementos'].append({
                'runs': [],
                'nivel': nivel
            })

        # Texto encontrado
        elif texto_run != "":
            # Añadir run a lista "runs"
            slide_list[-1]['placeholders'][-1]['elementos'][-1]['runs'].append(
                {
                    'texto': texto_run,
                    'negrita': negrita,
                    'cursiva': cursiva,
                    'subrayado': subrayado,
                    'tamaño_txt': tamaño_texto,
                    'color_txt': [r, g, b] if r != None else None,
                    'fuente': fuente
                })

        # Imagen encontrada
        elif img_path != None:
            # Añadir la imagen a los elementos del último placeholder
            slide_list[-1]['placeholders'][-1]['elementos'].append({
                'img_path': img_path,
                'img_param_dict': {
                    'img_modo': img_modo,
                    'rec_offset_izq_porc': rec_offset_izq_porc,
                    'rec_offset_der_porc': rec_offset_der_porc,
                    'rec_offset_sup_porc': rec_offset_sup_porc,
                    'rec_offset_inf_porc': rec_offset_inf_porc,
                    'ancho_img_cm': ancho_img_cm,
                    'alto_img_cm': alto_img_cm,
                    'img_alinear_v': img_alinear_v,
                    'img_alinear_h': img_alinear_h,
                }
            })

        # Tabla encontrada
        elif tabla != None:
            # Leer la hoja de Excel
            df_tabla = pd.read_excel(INPUT_EXCEL_CFG, sheet_name=tabla)

            # Extraer los formatos avanzados de los títulos
            formatos, df_tabla = ExtraerFormatosAvanzados(df_tabla)

            # Crear la lista para almacenar la tabla formateada
            tabla_format = []

            # Añadir los títulos sin formato a la primera fila
            tabla_format.append(list(df_tabla.columns))

            # Iterar sobre los valores del DataFrame
            for row_idx in range(df_tabla.shape[0]):
                row_values = []
                for col_idx in range(df_tabla.shape[1]):
                    # Obtener el título de la columna original (con formato=) para aplicar el formato
                    titulo_original = list(formatos.keys())[col_idx]
                    formato = formatos[titulo_original] # Obtener el formato original

                    # Obtener el valor del DataFrame
                    valor = df_tabla.iloc[row_idx, col_idx]

                    # Aplicar el formato avanzado basado en el título original
                    valor_formateado = AplicarFormatoAvanzado(valor, formato)

                    # Añadir el valor a la fila
                    row_values.append(valor_formateado)

                # Añadir la fila a la tabla
                tabla_format.append(row_values)

            # Añadir la tabla a los elementos del último placeholder
            slide_list[-1]['placeholders'][-1]['elementos'].append({
                'tabla': tabla_format,
                'estilo_tabla': estilo_tabla,
                'tamaño_txt': tamaño_texto,
                'fuente': fuente,
            })

    return slide_list


# Rutina principal
def main():
    # Eliminar y crear OUTPUT_PATH de nuevo
    RestaurarCarpeta(OUTPUT_PATH)

    # Crear Presentación
    prs = Presentation(INPUT_PPTX_TEMPLATE)

    # Imprimir los diseños de la presentación por consola
    ImprimirDiseñosPresentacion(prs)

    # Exportar el nombre de los placeholders en una nueva presentación
    ExportarNombrePlaceholders()

    # Importar datos desde fichero Excel
    slide_list = ImportarDatosExcel(prs)

    # Iteración sobre la lista de diccionarios SLIDE_LIST
    for slide_dict in slide_list:
        # Crear diapositiva
        layout_id = slide_dict['layout']
        slide = CrearDiapositiva(prs, layout_id)

        # Iterar sobre los placeholders
        for ph_dict in slide_dict['placeholders']:
            placeholder_id = ph_dict['placeholder_id']

            # Iterar sobre cada elemento del placeholder
            for elem_idx, elem_dict in enumerate(ph_dict['elementos']):
                # Texto detectado
                if ('nivel' in elem_dict):
                    # Añadir texto al placeholder
                    nivel = elem_dict['nivel']
                    runs_list = elem_dict['runs']
                    AñadirTextoPlaceholder(slide, placeholder_id, runs_list, nivel, elem_idx)

                # Imagen detectada
                elif ('img_path' in elem_dict):
                    ph = slide.placeholders[placeholder_id]
                    AñadirImagenPlaceholder(slide, elem_dict['img_path'], ph, elem_dict['img_param_dict'])

                # Tabla detectada
                elif ('tabla' in elem_dict):
                    AñadirTablaPlaceholder(slide, placeholder_id, elem_dict)

    # Guardar la presentación
    prs.save(OUTPUT_PATH + "\\" + "PlantillasProyecto_Output.pptx")

    # Mensaje de presentación creada
    print("")
    print("La presentación se ha creado correctamente!")


if __name__ == '__main__':
    main()

